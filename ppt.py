"""
# To run:
#   python ppt.py <path> <customer folder> <template ppt>
# or
#   python ppt.py
#
# Take a list of picture files (i.e. <path> + <customer_folder>)
# Find out the slide number for a given picture based on Identify the slide for the picture
#
"""

from datetime import date
import glob
import click
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import yaml
import logging

logging.basicConfig(level=logging.INFO)

def add_textbox_formatted(preso, slide_num, text_config, text):
    """ add a text box to slide with formatted text"""
    slide = preso.slides[slide_num]
    color = text_config['color']
    pos = text_config['pos']
    # left, top, width, height
    text_box = slide.shapes.add_textbox(Inches(pos[0]),
                                        Inches(pos[1]),
                                        Inches(pos[2]),
                                        Inches(pos[3]))
    text_frame = text_box.text_frame

    paragraph = text_frame.add_paragraph()
    run = paragraph.add_run()
    run.text = text
    font = run.font
    font.name = text_config['font']
    font.size = Pt(text_config['size'])
    paragraph.font.color.rgb = RGBColor(*color)

def find_size(slide_num, slide_types):
    """ find picture size based on size slide map """
    for slide_type in slide_types:
        if slide_num in slide_type['slides']:
            return slide_type['size']

    return None

def find_slide_num(filename, file_slide_map):
    """ find slide number based on file slide map """
    for i in file_slide_map:
        # removimg common suffix '_tsc.png' in png filenamess 
        filename_cleaned = filename[:-8]
        if filename_cleaned.endswith(i):
            return file_slide_map[i]

    # png file isn't mapped to any slide in the template. Skip
    return None


@click.command()
@click.option('-c', '--customer_folder', default='h_b_zachry_company_2019_11_18',
              help='Customer folder for png files')
@click.option('-p', '--plot_path', 
              default='/users/hdpcompass/development/doojung/compass-lite-v2/compasslite_tableau_automation_extract/png/',
              #default='plots/',
              help='Path for all plot png files')
@click.option('-t', '--template', default='templates/Prod_Template_NewLite.pptx',
              help='Template ppt path')
def create_ppt(plot_path, customer_folder, template):
    """ main function to create ppt """
    logging.info('=================================================================================')
    logging.info("plot_path: {} customer: {} template: {}".format(plot_path, customer_folder, template))
    logging.info('=================================================================================')

    # getting config file based on the template
    with open(f'config/{template[10:-4]}yaml') as f:
        config = yaml.load(f, yaml.FullLoader)

    slide_types = config['slide_types']
    file_slide_map = config['file_slide_map']
    folder = customer_folder
    customer = customer_folder[:-11]

    # customer name is in customer-2019-xx-xx format
    preso = Presentation(template)

    # Add customer and date to first slide
    title_config = config['title_text']
    date_config = config['date_text']
    add_textbox_formatted(preso, 0, title_config, customer.upper().replace("_", " "))
    today = date.today().strftime("%B %d, %Y")
    today1 = date.today().strftime("%Y%m")
    add_textbox_formatted(preso, 0, date_config, today)

    # Add plots to their respective slides
    for png in glob.glob(plot_path + folder + "/*.png"):
        logging.info(png.replace(plot_path + folder, ''))
        slide_num = find_slide_num(png, file_slide_map)
        if slide_num:
            slide = preso.slides[slide_num]
            size = find_size(slide_num, slide_types)
            if size:
                logging.info("png size: {}, slide_num: {}".format(size, slide_num))
                pic = slide.shapes.add_picture(png,
                                               Inches(size['left']),
                                               Inches(size['top']),
                                               Inches(size['width']),
                                               Inches(size['height']))
                slide.shapes._spTree.insert(2, pic._element)
            else:
                logging.warning("Slide {} is missing in the slide_types, check config file")

    preso.save("/users/hdpcompass/production/ppt_files/" + "RAP_" + customer.upper() + "_" + today1 + ".pptx")


if __name__ == '__main__':
    create_ppt()
